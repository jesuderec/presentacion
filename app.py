import streamlit as st
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
# Importamos MSO_SHAPE_TYPE para manejar formas si es necesario
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
import requests
import json
import os
import docx
from pypdf import PdfReader
import io
import re
import openai

# --- Configuración básica de registro ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# --- Configuración de la API ---
def get_api_key(model_name):
    if model_name == "deepseek-chat":
        return os.getenv("DEEPSEEK_API_KEY")
    elif "gpt" in model_name:
        return os.getenv("OPENAI_API_KEY")
    return None

def setup_openai_client(api_key):
    openai.api_key = api_key

# --- Optimización de texto ---
def optimize_text_for_ai(text_content):
    cleaned_text = re.sub(r'[^\w\s.,?!¡¿]', '', text_content, flags=re.UNICODE)
    optimized_text = re.sub(r'\s+', ' ', cleaned_text).strip()
    return optimized_text

# --- Generación de slides con la IA (Prompt mejorado estilo Gems) ---
def generate_slides_data_with_ai(texto_contenido_principal, texto_estructura_base, num_slides, model_name, api_key):
    texto_contenido_principal = optimize_text_for_ai(texto_contenido_principal)
    texto_estructura_base = optimize_text_for_ai(texto_estructura_base)

    prompt = f"""
    ## ROL Y OBJETIVO
    Actúa como un Docente Universitario experto. Tu misión es transformar el "DOCUMENTO FUENTE" en una presentación académica de alto impacto, enfocada exclusivamente en contenido textual de excelencia.

    ## CONTEXTO
    - **CONTENIDO PRINCIPAL:** "{texto_contenido_principal}"
    - **ESTRUCTURA GUÍA (Opcional):** "{texto_estructura_base}"

    ## INSTRUCCIONES DE CONTENIDO (ESTILO GEMS)
    1. **Fidelidad Académica:** Extrae la teoría y conceptos clave del documento fuente. El tono debe ser formal pero cercano (estilo aula universitaria).
    2. **Estructura:** Genera exactamente {num_slides + 2} diapositivas (incluyendo una Introducción y una Conclusión).
    3. **Regla de Oro de los Bullets:** Máximo 3 o 4 puntos por diapositiva. Cada punto debe tener un MÁXIMO de 20 palabras. Evita párrafos en los bullets; usa ideas fuerza.
    4. **Narrativa (El Guion):** Este es el elemento más importante. Debe ser un texto extenso y fluido (150-300 palabras) por diapositiva. Debe explicar el "porqué", usar analogías y ejemplos didácticos. No debe limitarse a leer las viñetas.

    ## FORMATO DE SALIDA (JSON ÚNICAMENTE)
    Tu respuesta debe ser un objeto JSON válido, con una clave raíz "slides", sin texto adicional.
    {{
      "slides": [
        {{
          "title": "Título Académico descriptivo",
          "bullets": ["Punto clave 1", "Punto clave 2", "Punto clave 3"],
          "narrative": "Narrativa completa y detallada para el docente..."
        }}
      ]
    }}
    """

    try:
        headers = {'Content-Type': 'application/json', 'Authorization': f'Bearer {api_key}'}
        ai_response_content = ""
        
        if "deepseek" in model_name:
            api_url = "https://api.deepseek.com/v1/chat/completions"
            payload = {
                "model": "deepseek-chat", 
                "messages": [{"role": "user", "content": prompt}], 
                "temperature": 0.2, 
                "response_format": {"type": "json_object"}
            }
            response = requests.post(api_url, headers=headers, data=json.dumps(payload))
            response.raise_for_status()
            ai_response_content = response.json()["choices"][0]["message"]["content"]
            
        elif "gpt" in model_name:
            setup_openai_client(api_key)
            response = openai.chat.completions.create(
                model="gpt-4o-mini", 
                messages=[{"role": "user", "content": prompt}], 
                response_format={"type": "json_object"}
            )
            ai_response_content = response.choices[0].message.content

        # Limpieza de JSON
        match = re.search(r'```(?:json)?\s*([\s\S]*?)\s*```', ai_response_content)
        clean_json_str = match.group(1) if match else ai_response_content
        
        parsed_data = json.loads(clean_json_str)
        return parsed_data
    except Exception as e:
        logging.error(f"Error al procesar con IA: {e}")
        return None

# --- Funciones para crear presentación (Ajustada para Texto Completo) ---
def create_presentation(slides_data, presentation_title, presentation_subtitle):
    try:
        prs = Presentation()
        # Colores institucionales (Ejemplo: Granate y Blanco)
        color_fondo = RGBColor(82, 0, 41)
        color_texto = RGBColor(255, 255, 255)
        color_subtitulo = RGBColor(200, 200, 200)

        # Configurar fondo del Master
        master = prs.slide_masters[0]
        fill = master.background.fill
        fill.solid()
        fill.fore_color.rgb = color_fondo

        title_slide_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]
        
        # 1. Slide de Título
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = presentation_title
        subtitle.text = presentation_subtitle
        
        # 2. Slides de Contenido
        for slide_info in slides_data.get("slides", []):
            slide = prs.slides.add_slide(content_layout)
            
            # Título
            title_shape = slide.shapes.title
            title_shape.text = slide_info.get("title", "")
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            title_shape.text_frame.paragraphs[0].font.color.rgb = color_texto

            # Cuerpo (Ampliado ya que no hay imagen)
            body_shape = slide.placeholders[1]
            body_shape.width = Inches(8.5) # Casi todo el ancho
            tf = body_shape.text_frame
            tf.clear() 
            
            # Bullets Estilo Gems
            for bullet_point in slide_info.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet_point
                p.font.color.rgb = color_texto
                p.font.size = Pt(22)
                p.level = 0
                p.space_after = Pt(12)
            
            # Narrativa integrada (Resumen visual)
            narrative_text = slide_info.get("narrative", "")
            if narrative_text:
                p_narr = tf.add_paragraph()
                # Mostramos un resumen del guion en itálica en la parte inferior
                p_narr.text = f"\nGuion Docente: {narrative_text[:200]}..."
                p_narr.font.size = Pt(13)
                p_narr.font.italic = True
                p_narr.font.color.rgb = color_subtitulo

        # 3. Slide Final
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "¡Muchas Gracias!"
        slide.placeholders[1].text = "Fin de la Sesión"

        return prs
    except Exception as e:
        logging.error(f"Error al crear la presentación: {e}")
        return None

# --- Funciones para leer archivos ---
def read_text_from_file(uploaded_file):
    if uploaded_file is None:
        return ""
    uploaded_file.seek(0)
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    if file_extension == ".txt":
        return uploaded_file.read().decode("utf-8")
    elif file_extension == ".pdf":
        reader = PdfReader(uploaded_file)
        text = ""
        for page in reader.pages:
            if page.extract_text():
                text += page.extract_text()
        return text
    elif file_extension == ".docx":
        doc = docx.Document(uploaded_file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    return ""

# --- Interfaz de Streamlit ---
st.set_page_config(page_title="Docente IA - Generador Pro", page_icon="🎓")
st.title("Generador de Presentaciones Académicas 🎓✨")
st.markdown("Crea presentaciones de nivel universitario con guiones detallados (Estilo Gems).")

with st.sidebar:
    st.header("⚙️ Configuración")
    model_text_option = st.selectbox("Modelo de Lenguaje:", ["gpt-4o-mini", "deepseek-chat"])
    max_text_length = st.slider("Límite de caracteres del documento:", 500, 15000, 6000)

st.header("📄 Información General")
presentation_title = st.text_input("Título de la Clase:", "Ej. Fundamentos de Macroeconomía")
presentation_subtitle = st.text_input("Subtítulo o Módulo:", "Ej. Unidad 1: El Ciclo Económico")
num_slides = st.slider("Número de diapositivas de desarrollo:", 3, 20, 7)

st.header("⚙️ Contenido Fuente")
uploaded_file_content = st.file_uploader("Sube el material de estudio", type=["txt", "docx", "pdf"])
text_input_content = st.text_area("O pega el texto aquí", height=200)

content_to_process = read_text_from_file(uploaded_file_content) if uploaded_file_content else text_input_content

is_button_disabled = not bool(presentation_title.strip() and content_to_process.strip())

col1, col2 = st.columns(2)

if col1.button("Generar Material Docente", disabled=is_button_disabled):
    content_truncated = content_to_process[:max_text_length]
    with st.spinner("El Docente IA está analizando y redactando..."):
        selected_ai_key = get_api_key(model_text_option)
        
        if not selected_ai_key:
            st.error("Error: No se encontró la API Key en las variables de entorno.")
        else:
            slides_data = generate_slides_data_with_ai(content_truncated, "", num_slides, model_text_option, selected_ai_key)
            
            if slides_data:
                # Crear PPTX
                prs = create_presentation(slides_data, presentation_title, presentation_subtitle)
                
                if prs:
                    # Guardar PPTX en memoria
                    pptx_file = BytesIO()
                    prs.save(pptx_file)
                    st.session_state.presentation_data = pptx_file.getvalue()
                    
                    # Generar Narrativa Completa para TXT
                    narrative_full_text = f"GUION DE CLASE: {presentation_title}\n" + "="*40 + "\n\n"
                    for i, slide in enumerate(slides_data.get("slides", [])):
                        narrative_full_text += f"DIAPOSITIVA {i+1}: {slide.get('title', '')}\n"
                        narrative_full_text += f"CONTENIDO VISUAL: {', '.join(slide.get('bullets', []))}\n\n"
                        narrative_full_text += f"NARRATIVA DEL DOCENTE:\n{slide.get('narrative', '')}\n"
                        narrative_full_text += "\n" + "-"*30 + "\n\n"
                    
                    st.session_state.narrative_data = narrative_full_text.encode('utf-8')
                    st.success("¡Material generado con éxito!")

if col2.button("Limpiar"):
    for key in ['presentation_data', 'narrative_data']:
        if key in st.session_state: del st.session_state[key]
    st.rerun()

# --- Zona de Descarga ---
if 'presentation_data' in st.session_state:
    st.markdown("---")
    st.header("✅ Descargas Disponibles")
    
    with st.expander("📝 Vista Previa del Guion Docente"):
        st.text(st.session_state.narrative_data.decode('utf-8'))
        
    st.download_button(
        label="📥 Descargar Presentación (PPTX)",
        data=st.session_state.presentation_data,
        file_name=f"Clase_{presentation_title.replace(' ', '_')}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    
    st.download_button(
        label="📥 Descargar Guion Completo (TXT)",
        data=st.session_state.narrative_data,
        file_name=f"Guion_{presentation_title.replace(' ', '_')}.txt",
        mime="text/plain"
    )
